import os
import asyncio
import requests
from telegram import Update, InlineKeyboardMarkup, InlineKeyboardButton
from telegram.ext import (
    Application,
    CommandHandler,
    MessageHandler,
    CallbackQueryHandler,
    ContextTypes,
    filters,
)
from dotenv import load_dotenv
from datetime import datetime, timedelta
from deep_translator import GoogleTranslator
from apscheduler.schedulers.asyncio import AsyncIOScheduler
import speech_recognition as sr
from gtts import gTTS
from pydub import AudioSegment
import io
import logging
from PIL import Image
from pptx import Presentation
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet

# Logging sozlamalari
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# .env yuklash
load_dotenv()
TELEGRAM_TOKEN = os.getenv("TELEGRAM_TOKEN")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
WEATHER_API_KEY = os.getenv("WEATHER_API_KEY")
ADMIN_ID = int(os.getenv("ADMIN_ID", "0"))
RENDER_EXTERNAL_HOSTNAME = os.getenv("RENDER_EXTERNAL_HOSTNAME")

if not TELEGRAM_TOKEN or not OPENAI_API_KEY or not WEATHER_API_KEY:
    logger.error("TELEGRAM_TOKEN, OPENAI_API_KEY yoki WEATHER_API_KEY topilmadi!")
    exit(1)

# OpenAI API sozlamalari
OPENAI_API_BASE_URL = "https://api.openai.com/v1/chat/completions"
OPENAI_MODEL = "gpt-3.5-turbo"

# OpenAI orqali matn generatsiya qilish funksiyasi
async def openai_generate_content(prompt, max_tokens=500, temperature=0.7, retries=3, delay=5):
    headers = {
        "Authorization": f"Bearer {OPENAI_API_KEY}",
        "Content-Type": "application/json",
        "User-Agent": "TelegramBot/1.0"
    }
    data = {
        "model": OPENAI_MODEL,
        "messages": [{"role": "user", "content": prompt}],
        "max_tokens": max_tokens,
        "temperature": temperature
    }
    for attempt in range(retries):
        try:
            logger.info(f"OpenAI API so‘rovi ({attempt + 1}/{retries}): {prompt[:50]}...")
            response = requests.post(OPENAI_API_BASE_URL, headers=headers, json=data, timeout=15)
            response.raise_for_status()
            logger.info("OpenAI API javobi muvaffaqiyatli olingan.")
            return response.json()["choices"][0]["message"]["content"]
        except requests.exceptions.RequestException as e:
            logger.error(f"OpenAI API xatoligi: {str(e)}, Status Code: {getattr(e.response, 'status_code', 'N/A')}")
            if getattr(e.response, 'status_code', None) == 429 and attempt < retries - 1:
                logger.info(f"429 xatoligi aniqlandi. {delay} soniya kutish va qayta urinish...")
                await asyncio.sleep(delay)
                delay *= 2  # Eksponensial orqa qadam
            else:
                if getattr(e.response, 'status_code', None) == 429:
                    return "Xatolik: Juda ko‘p so‘rov. Iltimos, biroz kuting yoki obunani ko‘rib chiqing."
                return f"Xatolik: {str(e)}"
    return "Xatolik: Maksimal qayta urinishlar soniga yetildi."

# 🌤 O‘zbekiston shaharlar ro‘yxati
UZ_CITIES = [
    "Tashkent", "Samarkand", "Bukhara", "Khiva", "Andijan", "Namangan",
    "Fergana", "Kokand", "Jizzakh", "Navoiy", "Qarshi", "Termez",
    "Gulistan", "Shahrisabz", "Urgench"
]

# 💰 Mashhur kriptovalyutalar
CRYPTO_COINS = ["bitcoin", "ethereum", "tether", "bnb", "solana", "dogecoin"]

# 🔤 Tarjima tillari
LANG_CODES = {
    "🇺🇸 Ingliz": "en",
    "🇷🇺 Rus": "ru",
    "🇺🇿 O‘zbek": "uz",
    "🇹🇷 Turk": "tr",
    "🇩🇪 Nemis": "de",
    "🇫🇷 Fransuz": "fr"
}

# 🌤 Inglizcha → O‘zbekcha ob-havo tarjimalari
WEATHER_CONDITIONS = {
    "clear sky": "ochiq osmon",
    "few clouds": "biroz bulutli",
    "scattered clouds": "sochma bulutlar",
    "broken clouds": "qisman bulutli",
    "overcast clouds": "to‘liq bulutli",
    "shower rain": "jala",
    "light rain": "yengil yomg‘ir",
    "moderate rain": "o‘rtacha yomg‘ir",
    "heavy intensity rain": "kuchli yomg‘ir",
    "rain": "yomg‘ir",
    "snow": "qor",
    "mist": "tuman",
    "thunderstorm": "momaqaldiroq",
    "fog": "tuman",
    "haze": "xira havo",
    "dust": "chang",
    "sand": "qumli bo‘ron",
    "tornado": "tornado"
}

# --- Foydalanuvchilarni saqlash ---
def add_user(user_id, context):
    users = context.bot_data.get("users", set())
    users.add(user_id)
    context.bot_data["users"] = users

# ---- Chat tarixini boshqarish ----
def get_chat_history(context, user_id, max_length=5):
    history = context.user_data.get(user_id, {}).get("chat_history", [])
    return history[-max_length:]

def update_chat_history(context, user_id, message):
    if user_id not in context.user_data:
        context.user_data[user_id] = {"chat_history": [], "language": "uz"}
    context.user_data[user_id]["chat_history"].append(message)

# ---- Ovozli xabarni matnga aylantirish ----
def speech_to_text(voice_file):
    recognizer = sr.Recognizer()
    try:
        with io.BytesIO(voice_file) as audio_file:
            audio = AudioSegment.from_file(audio_file, format="ogg")
            audio.export("temp.wav", format="wav")
            with sr.AudioFile("temp.wav") as source:
                audio_data = recognizer.record(source)
                text = recognizer.recognize_google(audio_data, language="uz-UZ")
                return text
    except sr.UnknownValueError:
        return "Ovozli xabar tushunilmadi."
    except sr.RequestError as e:
        return f"Audioni qayta ishlashda xatolik: {str(e)}"
    except Exception as e:
        return f"Xatolik: {str(e)}"
    finally:
        if os.path.exists("temp.wav"):
            os.remove("temp.wav")

# ---- Matnni ovozga aylantirish ----
def text_to_speech(text, lang="uz"):
    if lang == "uz":
        lang = "ru"  # gTTS uchun O‘zbek tilida yo‘q, shuning uchun Rus tiliga o‘taman
    try:
        tts = gTTS(text=text, lang=lang, slow=False)
        with io.BytesIO() as audio_file:
            tts.write_to_fp(audio_file)
            audio_file.seek(0)
            return audio_file
    except Exception as e:
        logger.error(f"Ovozga aylantirishda xatolik: {str(e)}")
        return None

# ---- Tilni o'zgartirish ----
async def set_language(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    keyboard = [[InlineKeyboardButton(name, callback_data=f"set_lang_{code}")]
                for name, code in LANG_CODES.items()]
    reply_markup = InlineKeyboardMarkup(keyboard)
    await update.message.reply_text("🔤 Qaysi tilni tanlaysiz?", reply_markup=reply_markup)

async def set_language_button(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    lang = query.data.replace("set_lang_", "")
    user_id = query.from_user.id
    context.user_data[user_id]["language"] = lang
    await query.edit_message_text(f"Til o'zgartirildi: {lang}")

# ---- Rasmlarni analiz qilish ----
async def handle_photo_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    photo = update.message.photo[-1]
    file = await context.bot.get_file(photo.file_id)
    image_data = await file.download_as_bytearray()
    logger.info("Rasm yuklab olindi.")

    try:
        base64_image = image_data.hex()  # Hex formatida aylantirish
        headers = {
            "Authorization": f"Bearer {OPENAI_API_KEY}",
            "Content-Type": "application/json",
            "User-Agent": "TelegramBot/1.0"
        }
        prompt = "Bu rasmni tahlil qiling va tavsiflang."
        data = {
            "model": "gpt-4o",  # OpenAI'da rasm tahlili uchun mos model
            "messages": [{"role": "user", "content": [
                {"type": "text", "text": prompt},
                {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
            ]}],
            "max_tokens": 500
        }
        response = requests.post(OPENAI_API_BASE_URL, headers=headers, json=data, timeout=15)
        response.raise_for_status()
        await update.message.reply_text(response.json()["choices"][0]["message"]["content"])
    except Exception as e:
        logger.error(f"Rasm analizida xatolik: {str(e)}")
        await update.message.reply_text(f"Rasm analizida xatolik: {str(e)}")

# ---- Umumiy xabar handleri (matn va ovoz) ----
async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    add_user(user_id, context)

    last_message_time = context.user_data.get("last_message_time", None)
    if last_message_time and datetime.now() < last_message_time + timedelta(seconds=10):  # 5 soniyadan 10 soniyaga oshirildi
        await update.message.reply_text("⏳ Iltimos, 10 soniya kuting!")
        return
    context.user_data["last_message_time"] = datetime.now()

    history = get_chat_history(context, user_id)

    if update.message.text and "target_lang" not in context.user_data:
        text = update.message.text
        logger.info(f"Matnli xabar: {text}")
        update_chat_history(context, user_id, {"user": text, "bot": ""})
    elif update.message.voice:
        voice = update.message.voice
        file = await context.bot.get_file(voice.file_id)
        voice_data = await file.download_as_bytearray()
        logger.info("Ovozli xabar yuklab olindi.")
        text = speech_to_text(voice_data)
        await update.message.reply_text(f"Sizning ovozli xabaringiz: {text}")
        update_chat_history(context, user_id, {"user": text, "bot": ""})
    elif update.message.photo:
        await handle_photo_message(update, context)
        return
    else:
        await update.message.reply_text("Noma'lum xabar. Matn, ovoz yoki rasm yuboring.")
        return

    if "target_lang" not in context.user_data:
        user_id = update.effective_user.id
        lang = context.user_data.get(user_id, {}).get("language", "uz")
        prompt = f"Javobni {lang} tilida bering. Har qanday mavzuda yordam bering:\n"
        prompt += "\n".join([f"Foydalanuvchi: {msg['user']}\nBot: {msg['bot']}" for msg in history])
        prompt += f"\nFoydalanuvchi: {text}\nBot: "
        try:
            response_text = await openai_generate_content(prompt)
            update_chat_history(context, user_id, {"user": text, "bot": response_text})
            await update.message.reply_text(f"Javob: {response_text}")
        except Exception as e:
            logger.error(f"OpenAI javobida xatolik: {str(e)}")
            await update.message.reply_text(f"Xatolik: {str(e)}")

# ---- Start va Help ----
async def start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    add_user(update.effective_user.id, context)
    await update.message.reply_text("Salom! Men AI yordamchiman 🤖. Har qanday savol bilan yordam beraman yoki /help ni ishlat.")

async def help_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    text = (
        "/start - Botni ishga tushirish\n"
        "/help - Yordam\n"
        "/weather - Shaharni tanlab, ob-havo olish\n"
        "/crypto - Kripto tanlab, narxini olish\n"
        "/translate - Tilni tanlab, tarjima qilish\n"
        "/currency - Bugungi valyuta kurslari (CBU)\n"
        "/presentation - AI yordamida prezentatsiya tayyorlash\n"
        "🤖 Har qanday savol uchun matn, ovoz yoki rasm yuboring – ketma-ket suhbatda javob beraman!\n"
    )
    if update.effective_user.id == ADMIN_ID:
        text += (
            "\n--- 🛠 Admin komandalar ---\n"
            "/broadcast - Hammaga xabar yuborish\n"
            "/report - So‘rovlar haqida hisobot\n"
            "/myid - O‘z ID’ingizni bilish\n"
        )
    await update.message.reply_text(text)

# ---- My ID ----
async def myid(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text(f"Sizning ID: {update.effective_user.id}")

# ---- WEATHER ----
async def weather_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    keyboard = [[InlineKeyboardButton(city, callback_data=f"weather_{city}")] for city in UZ_CITIES]
    reply_markup = InlineKeyboardMarkup(keyboard)
    await update.message.reply_text("🌤 Qaysi shahar ob-havosini bilmoqchisiz?", reply_markup=reply_markup)

async def weather_button(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    city = query.data.replace("weather_", "")
    url = f"http://api.openweathermap.org/data/2.5/weather?q={city}&appid={WEATHER_API_KEY}&units=metric&lang=en"
    try:
        res = requests.get(url, timeout=10).json()
        if res.get("cod") != 200:
            await query.edit_message_text(f"❌ Ob-havo topilmadi: {city}")
            return
        temp = res["main"]["temp"]
        desc = res["weather"][0]["description"].lower()
        uz_desc = WEATHER_CONDITIONS.get(desc, desc)
        await query.edit_message_text(f"🌤 {city} ob-havosi:\n{temp}°C, {uz_desc}")
    except Exception as e:
        await query.edit_message_text(f"Xatolik: {str(e)}")

# ---- CRYPTO ----
async def crypto_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    keyboard = [[InlineKeyboardButton(coin.capitalize(), callback_data=f"crypto_{coin}")]
                for coin in CRYPTO_COINS]
    reply_markup = InlineKeyboardMarkup(keyboard)
    await update.message.reply_text("💰 Qaysi kripto narxini bilmoqchisiz?", reply_markup=reply_markup)

async def crypto_button(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    coin = query.data.replace("crypto_", "")
    url = f"https://api.coingecko.com/api/v3/simple/price?ids={coin}&vs_currencies=usd"
    try:
        res = requests.get(url, timeout=10).json()
        if coin not in res:
            await query.edit_message_text(f"❌ Kripto topilmadi: {coin}")
            return
        price = res[coin]["usd"]
        await query.edit_message_text(f"💰 {coin.capitalize()} narxi: ${price}")
    except Exception as e:
        await query.edit_message_text(f"Xatolik: {str(e)}")

# ---- TRANSLATE ----
async def translate_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    keyboard = [[InlineKeyboardButton(name, callback_data=f"lang_{code}")]
                for name, code in LANG_CODES.items()]
    reply_markup = InlineKeyboardMarkup(keyboard)
    await update.message.reply_text("🔤 Qaysi tilga tarjima qilmoqchisiz?", reply_markup=reply_markup)

async def lang_button(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    lang = query.data.replace("lang_", "")
    context.user_data["target_lang"] = lang
    await query.edit_message_text(f"✍️ Endi matn yuboring, men uni `{lang}` tiliga tarjima qilaman.")

async def translate_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if "target_lang" in context.user_data:
        lang = context.user_data["target_lang"]
        text = update.message.text
        try:
            prompt = f"Translate the following text into {lang}: {text}"
            response_text = await openai_generate_content(prompt)
            await update.message.reply_text(f"🔤 Tarjima ({lang}): {response_text}")
            del context.user_data["target_lang"]
        except Exception as e:
            logger.error(f"Tarjima xatoligi: {str(e)}")
            try:
                translated = GoogleTranslator(source="auto", target=lang).translate(text)
                await update.message.reply_text(f"🔤 Tarjima ({lang}): {translated}")
                del context.user_data["target_lang"]
            except Exception as e2:
                await update.message.reply_text(f"❌ Tarjima xatoligi: {str(e2)}")

# ---- CURRENCY ----
async def currency(update: Update, context: ContextTypes.DEFAULT_TYPE):
    url = "https://cbu.uz/oz/arkhiv-kursov-valyut/json/"
    try:
        res = requests.get(url, timeout=10).json()
        if not res:
            await update.message.reply_text("❌ Valyuta kurslari topilmadi.")
            return
        selected = [c for c in res if c["Ccy"] in ["USD", "EUR", "RUB"]]
        text = "💱 Bugungi valyuta kurslari (CBU):\n\n"
        for c in selected:
            text += f"1 {c['Ccy']} = {c['Rate']} so‘m\n"
        await update.message.reply_text(text)
    except Exception as e:
        await update.message.reply_text(f"Xatolik: {str(e)}")

# ---- ADMIN funksiyalari ----
async def broadcast(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("❌ Siz admin emassiz.")
        return
    if not context.args:
        await update.message.reply_text("Foydalanish: /broadcast Xabar matni")
        return
    text = " ".join(context.args)
    users = context.bot_data.get("users", set())
    if not users:
        await update.message.reply_text("📭 Hali foydalanuvchi yo‘q.")
        return
    sent, failed = 0, 0
    for uid in users:
        try:
            await context.bot.send_message(uid, f"📢 Admin xabari:\n\n{text}")
            sent += 1
        except Exception:
            failed += 1
    await update.message.reply_text(f"✅ Yuborildi: {sent} ta\n❌ Xato: {failed} ta")

async def report(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ADMIN_ID:
        await update.message.reply_text("❌ Siz admin emassiz.")
        return
    await send_report(context)

async def send_report(context: ContextTypes.DEFAULT_TYPE):
    logs = context.bot_data.get("logs", [])
    users = context.bot_data.get("users", set())
    if not logs:
        await context.bot.send_message(ADMIN_ID, "📊 Bugun hech qanday so‘rov bo‘lmadi.")
        return
    msg = (
        f"📊 Kunlik hisobot\n\n"
        f"👥 Foydalanuvchilar soni: {len(users)}\n"
        f"💬 So‘rovlar soni: {len(logs)}\n\n"
        "📝 Oxirgi 5 ta so‘rov:\n"
    )
    for time, uid, text in logs[-5:]:
        msg += f"🕒 {time} | 👤 {uid}\n💬 {text}\n\n"
    await context.bot.send_message(ADMIN_ID, msg)

# ---- Prezentatsiya tayyorlash ----
async def presentation_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.message.reply_text("🎥 Prezentatsiya mavzusini kiriting (masalan, 'O‘zbekiston tarixi' yoki 'AI texnologiyalari').")
    context.user_data["awaiting_presentation_topic"] = True

async def handle_presentation_topic(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if context.user_data.get("awaiting_presentation_topic"):
        topic = update.message.text
        await update.message.reply_text(f"📝 '{topic}' bo'yicha prezentatsiya tayyorlanmoqda... Iltimos, kuting.")
        
        prompt = f"Siz AI yordamchisiz. Quyidagi mavzu bo'yicha qisqa prezentatsiya matni yarating: {topic}. Strukturani quyidagi tarzda saqlang:\n- Sarlavha\n- Kirish (2-3 jumlali)\n- Asosiy qism (3 ta asosiy nuqta bilan)\n- Xulosa (1-2 jumlali)\nNatijani faqat matn sifatida qaytaring, hech qanday qo'shimcha izohsiz."
        try:
            response_text = await openai_generate_content(prompt, max_tokens=1000)
            presentation_text = response_text.split('\n')
            logger.info(f"Generatsiya qilingan matn: {presentation_text}")

            ppt = Presentation()
            title_slide_layout = ppt.slide_layouts[0]
            slide = ppt.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = presentation_text[0] if presentation_text else topic
            subtitle.text = "Tayyorlandi: " + datetime.now().strftime("%Y-%m-%d %H:%M")

            body_slide_layout = ppt.slide_layouts[1]
            for line in presentation_text[1:]:
                if line.strip():
                    slide = ppt.slides.add_slide(body_slide_layout)
                    title = slide.shapes.title
                    body = slide.placeholders[1]
                    title.text = line.strip()
                    body.text = "Tafsilotlar bu yerda bo‘ladi..."

            ppt_io = io.BytesIO()
            ppt.save(ppt_io)
            ppt_io.seek(0)
            logger.info("PowerPoint fayli muvaffaqiyatli yaratildi.")

            pdf_io = io.BytesIO()
            doc = SimpleDocTemplate(pdf_io, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            for line in presentation_text:
                if line.strip():
                    story.append(Paragraph(line.strip(), styles['Heading1']))
                    story.append(Spacer(1, 12))
            doc.build(story)
            pdf_io.seek(0)
            logger.info("PDF fayli muvaffaqiyatli yaratildi.")

            await update.message.reply_document(document=ppt_io, filename=f"{topic}_presentation.pptx")
            await update.message.reply_document(document=pdf_io, filename=f"{topic}_presentation.pdf")
            await update.message.reply_text("🎉 Prezentatsiya PowerPoint (.pptx) va PDF formatida yuborildi!")
        except Exception as e:
            logger.error(f"Prezentatsiya yaratishda xatolik: {str(e)}")
            await update.message.reply_text(f"❌ Prezentatsiya yaratishda xatolik: {str(e)}")
        
        del context.user_data["awaiting_presentation_topic"]

# ---- MAIN ----
def main():
    application = Application.builder().token(TELEGRAM_TOKEN).build()

    # Buyruqlar
    application.add_handler(CommandHandler("start", start))
    application.add_handler(CommandHandler("help", help_command))
    application.add_handler(CommandHandler("weather", weather_start))
    application.add_handler(CommandHandler("crypto", crypto_start))
    application.add_handler(CommandHandler("translate", translate_start))
    application.add_handler(CommandHandler("currency", currency))
    application.add_handler(CommandHandler("broadcast", broadcast))
    application.add_handler(CommandHandler("report", report))
    application.add_handler(CommandHandler("myid", myid))
    application.add_handler(CommandHandler("presentation", presentation_start))

    # Tugma handlerlar
    application.add_handler(CallbackQueryHandler(weather_button, pattern="^weather_"))
    application.add_handler(CallbackQueryHandler(crypto_button, pattern="^crypto_"))
    application.add_handler(CallbackQueryHandler(lang_button, pattern="^lang_"))
    application.add_handler(CallbackQueryHandler(set_language_button, pattern="^set_lang_"))

    # Umumiy xabar handleri
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND | filters.VOICE | filters.PHOTO, handle_message))
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, translate_message))

    # Scheduler ni moslashtirish
    scheduler = AsyncIOScheduler()
    scheduler.add_job(send_report, "cron", hour=23, minute=59, args=[application])

    # Webhook yoki polling
    port = int(os.environ.get("PORT", 8443))
    url_path = TELEGRAM_TOKEN
    webhook_url = f"https://{RENDER_EXTERNAL_HOSTNAME}/{url_path}" if RENDER_EXTERNAL_HOSTNAME else None

    if not webhook_url:
        logger.info("Xato: RENDER_EXTERNAL_HOSTNAME aniqlanmadi! Polling rejimida ishlayapman.")
        application.run_polling(allowed_updates=Update.ALL_TYPES)
    else:
        logger.info(f"Bot webhook bilan ishga tushmoqda: {webhook_url}")
        application.run_webhook(
            listen="0.0.0.0",
            port=port,
            url_path=url_path,
            webhook_url=webhook_url,
            allowed_updates=Update.ALL_TYPES
        )

    # Scheduler ni botning loop'ida boshlash
    try:
        loop = asyncio.get_event_loop()
        if loop.is_running():
            loop.create_task(scheduler.start())
    except RuntimeError:
        logger.warning("Event loop allaqachon yopilgan yoki ishlamayapti.")

if __name__ == "__main__":
    main()